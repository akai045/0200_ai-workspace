# -*- coding: utf-8 -*-
"""PPTX レイアウト自動検査（はみ出し・重なり・テキスト溢れ検知）.

使い方:
    python tools/pptx_layout_check.py <file.pptx> [--json]

検査内容（geometry ベース・Office 不要・どのPCでも実行可）:
  1. off-slide      : 図形がスライド枠外へはみ出していないか
  2. overlap        : 画像とテキストボックスが重なっていないか
  3. text-overflow  : テキスト量がボックス高さを超える可能性（推定・警告）

終了コード: 問題なし=0 / 検出あり=1 / 実行エラー=2
ビルドに組み込む場合はこのファイルを import して check(path) を呼ぶ。
"""
import sys, json

def _in(emu):
    return emu / 914400.0 if emu is not None else 0.0

def check(path, char_pt_ratio=0.62, line_h_in=0.24):
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation(path)
    W = prs.slide_width; Hh = prs.slide_height
    tol = int(Inches(0.06))
    warns = []
    for si, slide in enumerate(prs.slides, 1):
        pics = []; boxes = []
        for sh in slide.shapes:
            L, T, Wd, Hg = sh.left, sh.top, sh.width, sh.height
            if None in (L, T, Wd, Hg):
                continue
            if L < -tol or T < -tol or L + Wd > W + tol or T + Hg > Hh + tol:
                warns.append({"slide": si, "type": "off-slide", "shape_id": sh.shape_id,
                              "detail": f"{sh.shape_type}"})
            if sh.shape_type == 13:  # PICTURE
                pics.append((L, T, Wd, Hg, sh))
            if sh.has_text_frame and sh.text_frame.text.strip():
                boxes.append((L, T, Wd, Hg, sh))
                usable_w = _in(Wd) - 0.3
                lines = 0
                for para in sh.text_frame.paragraphs:
                    t = para.text
                    if not t:
                        lines += 1; continue
                    sizes = [r.font.size.pt for r in para.runs if r.font.size] or [11]
                    sz = max(sizes)
                    cpl = max(1, int(usable_w * 72 / (sz * char_pt_ratio)))
                    lines += max(1, -(-len(t) // cpl))
                need = lines * line_h_in
                if need > _in(Hg) + 0.2:
                    warns.append({"slide": si, "type": "text-overflow", "shape_id": sh.shape_id,
                                  "detail": f"need~{need:.1f}in > box {_in(Hg):.1f}in"})
        for (l1, t1, w1, h1, p1) in pics:
            for (l2, t2, w2, h2, b2) in boxes:
                ix = max(0, min(l1 + w1, l2 + w2) - max(l1, l2))
                iy = max(0, min(t1 + h1, t2 + h2) - max(t1, t2))
                if ix > tol and iy > tol:
                    warns.append({"slide": si, "type": "overlap", "shape_id": p1.shape_id,
                                  "detail": f"picture {p1.shape_id} ↔ text {b2.shape_id}"})
    return {"file": path, "slides": len(prs.slides._sldIdLst), "warnings": warns}

def main():
    args = [a for a in sys.argv[1:] if not a.startswith("--")]
    as_json = "--json" in sys.argv
    if not args:
        print("usage: python tools/pptx_layout_check.py <file.pptx> [--json]"); sys.exit(2)
    try:
        res = check(args[0])
    except Exception as e:
        print("ERROR:", repr(e)); sys.exit(2)
    if as_json:
        print(json.dumps(res, ensure_ascii=False, indent=2))
    else:
        print(f"slides={res['slides']}  warnings={len(res['warnings'])}")
        for w in res["warnings"][:100]:
            print(f"  S{w['slide']}: {w['type']} — {w['detail']}")
        if not res["warnings"]:
            print("LAYOUT OK: はみ出し・重なり・テキスト溢れは検出されませんでした")
    sys.exit(1 if res["warnings"] else 0)

if __name__ == "__main__":
    main()
